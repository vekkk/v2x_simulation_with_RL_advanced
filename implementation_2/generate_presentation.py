import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

RESULTS_DIR = "data/results"
OUT_PPT = os.path.join(RESULTS_DIR, "Vehicular_DRL_Demo.pptx")

prs = Presentation()

# Helpers

def title_slide(title, subtitle=""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def bullet_slide(title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = b
        p.level = 0


def image_slide(title, image_paths):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = Inches(0.5)
    top = Inches(1.5)
    max_w = Inches(4.5)
    x = left
    for i, path in enumerate(image_paths):
        if os.path.exists(path):
            pic = slide.shapes.add_picture(path, x, top, width=max_w)
            x = x + max_w + Inches(0.2)

# Slides

title_slide("Hybrid RAT Selection with DQN", "Vehicular Network DRL Demo")

bullet_slide("Problem Statement", [
    "Vehicles must select the best Radio Access Technology (RAT) per timestep.",
    "Objectives: Maximize PRR & throughput, minimize latency & energy.",
    "Constraints: QoS thresholds (e.g., latency < 80 ms, PRR > 0.85).",
])

bullet_slide("DRL Solution", [
    "DQN (24-dim state → 128-128-64 → 4 actions).",
    "Action masking based on availability with fallback.",
    "Reward combines PRR, latency, throughput, energy, handover, QoS penalties.",
])

bullet_slide("System Architecture", [
    "Lightweight network simulator (RSSI, load, KPIs).",
    "Gym env wrapper for training/eval.",
    "Baselines: static, random, QoS-aware, signal, load-balance.",
])

# Results visuals
summary_csv = os.path.join(RESULTS_DIR, "summary_table.csv")
if os.path.exists(summary_csv):
    df = pd.read_csv(summary_csv)
    top = df.sort_values("prr", ascending=False)[["policy", "prr", "latency_ms", "throughput_mbps", "energy_j", "qos_rate"]]
    bullets = [f"{row.policy}: PRR={row.prr:.3f}, Lat={row.latency_ms:.1f} ms, Thr={row.throughput_mbps:.2f} Mbps, QoS={100*(1-row.qos_rate):.1f}%" for _, row in top.iterrows()]
    bullet_slide("Aggregate Results", bullets[:6])

image_slide("PRR / Latency / Throughput / Energy", [
    os.path.join(RESULTS_DIR, "bar_prr.png"),
    os.path.join(RESULTS_DIR, "bar_latency_ms.png"),
])
image_slide("QoS & Usage", [
    os.path.join(RESULTS_DIR, "bar_qos_rate.png"),
    os.path.join(RESULTS_DIR, "pie_usage_drl_dqn.png"),
])

bullet_slide("Contributions", [
    "End-to-end runnable demo (sim → DRL → baselines → eval → plots → PPT).",
    "Config-driven; easy to tune episodes/weights.",
    "Clear deltas vs. baselines and learning curve.",
])

bullet_slide("Future Work / IEEE Roadmap", [
    "Add realistic mobility traces & channel (3GPP/ETSI).",
    "Multi-vehicle, interference & scheduling models.",
    "PPO/SAC; multi-objective RL & constrained RL.",
])

prs.save(OUT_PPT)
print("Saved:", OUT_PPT)
